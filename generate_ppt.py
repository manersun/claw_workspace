#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    prs = Presentation()
    
    # 设置宽屏 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 品牌色
    GREEN = RGBColor(0x2D, 0xA0, 0x68)
    DARK_GREEN = RGBColor(0x1E, 0x6B, 0x45)
    LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xED)
    ORANGE = RGBColor(0xFF, 0x8C, 0x42)
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        
        return slide
    
    def add_content_slide(title, content_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, line in enumerate(content_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.space_after = Pt(12)
            if line.startswith("•"):
                p.level = 0
        
        return slide
    
    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # 标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN
        
        # 表格
        table_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers),
            Inches(0.5), Inches(1.5),
            Inches(12), Inches(5)
        )
        table = table_shape.table
        
        # 设置列宽
        for i in range(len(headers)):
            table.columns[i].width = Inches(12 / len(headers))
        
        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN
            tf = cell.text_frame
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 数据行
        for row_idx, row in enumerate(rows, 1):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)
                tf = cell.text_frame
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREEN
        
        return slide
    
    # ========== Slide 1: 封面 ==========
    slide = add_title_slide("邻选社区", "让邻里生活更美好\n社区团购新一代平台")
    subtitle = slide.placeholders[1]
    subtitle.text = "商业计划书\n\n[你的姓名] | 创始人 & CEO\n2026 年 3 月"
    
    # ========== Slide 2: 痛点与机会 ==========
    slide = add_content_slide("痛点与机会", [
        "• 消费者痛点：",
        "  - 买菜贵：传统菜场加价 30-50%",
        "  - 不新鲜：多层流通，到店已不新鲜",
        "  - 浪费时间：下班绕路买菜，排队结账",
        "",
        "• 市场机会：",
        "  - 社区团购市场规模：2025 年预计 3000 亿",
        "  - 渗透率仅 15%，仍有巨大增长空间",
        "  - 后疫情时代，线上买菜已成习惯",
        "  - 巨头收缩，区域玩家迎来机会窗口"
    ])
    
    # ========== Slide 3: 解决方案 ==========
    slide = add_content_slide("解决方案", [
        "• 邻选社区模式：",
        "  产地/批发商 → 城市仓 → 社区团长 → 消费者",
        "",
        "• 核心价值：",
        "  - 消费者：便宜 30% + 新鲜 + 送货上门",
        "  - 团长：佣金收入 10-15% + 零成本创业",
        "  - 供应商：稳定订单 + 降低损耗 + 快速回款",
        "",
        "• 差异化亮点：",
        "  - 精选 SKU：200 个高频爆品",
        "  - 品质管控：坏果包赔，24 小时售后",
        "  - 邻里社交：拼团 + 分享，增强粘性"
    ])
    
    # ========== Slide 4: 市场规模 ==========
    slide = add_content_slide("市场规模", [
        "• TAM (总可服务市场)：",
        "  - 中国生鲜零售市场：5.8 万亿/年",
        "  - 社区团购可渗透市场：8000 亿/年",
        "",
        "• SAM (可服务市场)：",
        "  - 目标城市（一线 + 新一线）：3000 亿/年",
        "  - 目标用户：25-45 岁家庭采购决策者 1.2 亿人",
        "",
        "• SOM (可获得市场)：",
        "  - 3 年目标：覆盖 10 个城市，50 亿 GMV/年",
        "  - 市占率目标：目标城市的 3-5%"
    ])
    
    # ========== Slide 5: 商业模式 ==========
    headers = ["项目", "占比"]
    rows = [
        ["商品成本", "75-80%"],
        ["团长佣金", "10-12%"],
        ["物流配送", "5-8%"],
        ["运营人力", "3-5%"],
        ["技术&其他", "2-3%"]
    ]
    slide = add_table_slide("商业模式 - 成本结构", headers, rows)
    
    # ========== Slide 6: 竞争优势 ==========
    headers = ["维度", "巨头", "邻选社区"]
    rows = [
        ["SKU 数量", "1000+", "200 精选"],
        ["品质管控", "标准化", "本地化严选"],
        ["团长管理", "弱管控", "强培训 + 激励"],
        ["售后响应", "48h", "24h"],
        ["决策速度", "慢", "快"],
        ["本地供应链", "弱", "强"]
    ]
    slide = add_table_slide("竞争优势 - vs 美团优选/多多买菜", headers, rows)
    
    # ========== Slide 7: 运营策略 ==========
    slide = add_content_slide("运营策略", [
        "• 冷启动计划（0-3 个月）：",
        "  - 选择 3 个标杆小区试点",
        "  - 招募 30 个种子团长（宝妈/便利店主）",
        "  - 打磨 50 个核心 SKU",
        "  - 验证单仓模型",
        "",
        "• 扩张节奏（4-12 个月）：",
        "  - 月新增 5-8 个社区",
        "  - 团长总数达 200 人",
        "  - 覆盖 50 个小区",
        "  - 日 GMV 突破 50 万"
    ])
    
    # ========== Slide 8: 财务预测 ==========
    headers = ["项目", "第 1 年", "第 2 年", "第 3 年"]
    rows = [
        ["GMV", "6,000 万", "30,000 万", "80,000 万"],
        ["营收", "1,200 万", "6,000 万", "16,000 万"],
        ["毛利", "240 万", "1,200 万", "3,200 万"],
        ["运营成本", "400 万", "1,000 万", "2,400 万"],
        ["净利润", "-160 万", "200 万", "800 万"]
    ]
    slide = add_table_slide("财务预测（三年）", headers, rows)
    
    # ========== Slide 9: 团队介绍 ==========
    slide = add_content_slide("团队介绍", [
        "• 创始人 & CEO - [你的名字]",
        "  [待补充：相关背景]",
        "  擅长：战略规划、资源整合",
        "",
        "• 联合创始人 & COO - [待招募]",
        "  目标背景：生鲜/零售运营经验",
        "  擅长：供应链管理、团队管理",
        "",
        "• 技术负责人 - [待招募]",
        "  目标背景：电商平台技术经验",
        "  擅长：系统架构、数据驱动",
        "",
        "💡 投资人提示：团队是早期投资的核心",
        "   建议尽快补齐运营和技术合伙人"
    ])
    
    # ========== Slide 10: 融资计划 ==========
    headers = ["项目", "详情"]
    rows = [
        ["融资金额", "300-500 万"],
        ["出让股份", "15-20%"],
        ["估值", "2000-2500 万（投前）"],
        ["资金用途", "供应链 40% / 团队 20% / 技术 15% / 市场 15% / 备用 10%"],
        ["退出路径", "3 年并购 / 5 年 IPO"]
    ]
    slide = add_table_slide("融资计划 - 天使轮", headers, rows)
    
    # ========== Slide 11: 里程碑 ==========
    headers = ["时间", "里程碑"]
    rows = [
        ["M1-3", "完成团队搭建，3 个试点小区运营"],
        ["M4-6", "验证单仓模型，扩展至 15 个小区"],
        ["M7-9", "覆盖 50 个小区，日 GMV 30 万+"],
        ["M10-12", "启动 Pre-A 轮融资，准备城市扩张"]
    ]
    slide = add_table_slide("里程碑（融资后 12 个月）", headers, rows)
    
    # ========== Slide 12: 联系我们 ==========
    slide = add_content_slide("联系我们", [
        "邻选社区",
        "",
        "📧 邮箱：[你的邮箱]",
        "📱 电话：[你的电话]",
        "📍 地址：[公司地址]",
        "",
        "感谢聆听！",
        "期待与您携手共创美好社区生活 🚀"
    ])
    
    # 保存
    prs.save('/Users/sunyu.maner/.openclaw/workspace/邻选社区 - 融资路演.pptx')
    print("✅ PPT 生成成功！")
    print("文件位置：/Users/sunyu.maner/.openclaw/workspace/邻选社区 - 融资路演.pptx")

if __name__ == "__main__":
    create_ppt()
