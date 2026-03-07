#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建 Prompt 设计指南 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 创建演示文稿
prs = Presentation()

# 设置宽屏 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_items):
    """添加内容页（项目符号）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        if item.startswith('  '):
            p.level = 1
            p.text = item.strip()
    
    return slide

def add_section_slide(prs, title):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 添加背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2),
        prs.slide_width, Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    shape.line.fill.background()
    
    # 添加标题
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        prs.slide_width - Inches(1), Inches(2)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    return slide

def add_code_slide(prs, title, code, explanation=None):
    """添加代码示例页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(0.8)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # 代码框
    codeBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        prs.slide_width - Inches(1), Inches(3.5)
    )
    codeBox.fill.solid()
    codeBox.fill.fore_color.rgb = RGBColor(40, 44, 52)
    codeBox.line.color.rgb = RGBColor(100, 100, 100)
    
    tf = codeBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = code
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 说明
    if explanation:
        expBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(5),
            prs.slide_width - Inches(1), Inches(2)
        )
        tf = expBox.text_frame
        p = tf.add_paragraph()
        p.text = explanation
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def add_comparison_slide(prs, title, bad_example, good_example):
    """添加对比示例页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        prs.slide_width - Inches(1), Inches(0.6)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 左侧：错误示例
    leftBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(1),
        Inches(6), Inches(5.5)
    )
    leftBox.fill.solid()
    leftBox.fill.fore_color.rgb = RGBColor(231, 76, 60)
    leftBox.line.fill.background()
    
    tf = leftBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "❌ 错误示例\n\n" + bad_example
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 右侧：正确示例
    rightBox = slide.shapes.add_textbox(
        Inches(6.8), Inches(1),
        Inches(6), Inches(5.5)
    )
    rightBox.fill.solid()
    rightBox.fill.fore_color.rgb = RGBColor(39, 174, 96)
    rightBox.line.fill.background()
    
    tf = rightBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "✅ 正确示例\n\n" + good_example
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    return slide

# ============ 开始创建幻灯片 ============

# 1. 标题页
add_title_slide(
    prs,
    "Prompt 设计指南",
    "如何设计高效的 AI 提示词\n\n" + datetime.now().strftime("%Y 年 %m 月")
)

# 2. 目录
add_content_slide(prs, "目录", [
    "1. 什么是 Prompt",
    "2. 为什么 Prompt 很重要",
    "3. Prompt 设计核心原则",
    "4. 常用 Prompt 框架",
    "5. 实战示例",
    "6. 常见错误与优化",
    "7. 高级技巧",
    "8. 总结与实践"
])

# 3. 章节：什么是 Prompt
add_section_slide(prs, "1. 什么是 Prompt")

# 4. Prompt 定义
add_content_slide(prs, "什么是 Prompt？", [
    "Prompt = 提示词 / 提示语",
    "",
    "定义：用户输入给 AI 的指令或问题",
    "",
    "作用：引导 AI 理解任务并生成期望的输出",
    "",
    "示例：",
    "  • 简单：\"写一首诗\"",
    "  • 复杂：\"你是一位资深产品经理，请为一款新的时间管理 APP 撰写产品需求文档，包含用户画像、核心功能、技术架构...\""
])

# 5. 为什么重要
add_content_slide(prs, "为什么 Prompt 很重要？", [
    "同样的 AI 模型，不同的 Prompt = 完全不同的结果",
    "",
    "好 Prompt 的价值：",
    "  • 输出质量提升 10 倍以上",
    "  • 减少反复修改的时间",
    "  • 降低 AI「幻觉」概率",
    "  • 让 AI 理解复杂任务",
    "",
    "一句话：Prompt 是与 AI 沟通的艺术"
])

# 6. 章节：核心原则
add_section_slide(prs, "2. Prompt 设计核心原则")

# 7. CLEAR 原则
add_content_slide(prs, "CLEAR 原则", [
    "C - Concise（简洁）",
    "  • 用词精准，避免冗余",
    "",
    "L - Logical（逻辑清晰）",
    "  • 结构化表达，分步骤说明",
    "",
    "E - Explicit（明确具体）",
    "  • 清晰定义任务、角色、输出格式",
    "",
    "A - Adaptive（适应性）",
    "  • 根据场景调整详细程度",
    "",
    "R - Reflective（可反思）",
    "  • 鼓励 AI 自我验证和优化"
])

# 8. 关键要素
add_content_slide(prs, "Prompt 的关键要素", [
    "1. 角色（Role）",
    "  • 「你是一位资深 Python 工程师」",
    "",
    "2. 任务（Task）",
    "  • 「请编写一个 Web 爬虫」",
    "",
    "3. 上下文（Context）",
    "  • 「用于监控竞争对手价格变化」",
    "",
    "4. 约束（Constraints）",
    "  • 「使用 requests 库，添加错误处理」",
    "",
    "5. 输出格式（Format）",
    "  • 「输出完整代码，包含注释」"
])

# 9. 章节：常用框架
add_section_slide(prs, "3. 常用 Prompt 框架")

# 10. RTF 框架
add_content_slide(prs, "RTF 框架（Role-Task-Format）", [
    "最简单的万能框架",
    "",
    "R - Role（角色）",
    "  • 你是一位经验丰富的数据分析师",
    "",
    "T - Task（任务）",
    "  • 分析这份销售数据，找出增长趋势",
    "",
    "F - Format（格式）",
    "  • 用表格展示，包含关键指标和洞察"
])

# 11. RTF 示例
add_code_slide(
    prs,
    "RTF 框架示例",
    """你是一位资深产品经理（Role）

请为一款新的健身 APP 设计用户增长方案（Task）

输出要求（Format）：
- 包含 3 个核心增长策略
- 每个策略说明实施步骤和预期效果
- 用 Markdown 表格总结""",
    "这个 Prompt 清晰定义了角色、任务和输出格式，AI 能准确理解需求"
)

# 12. CO-STAR 框架
add_content_slide(prs, "CO-STAR 框架（更完整）", [
    "C - Context（背景）",
    "  • 项目背景、目标用户、业务场景",
    "",
    "O - Objective（目标）",
    "  • 要解决什么问题，达到什么效果",
    "",
    "S - Style（风格）",
    "  • 专业、幽默、简洁、详细...",
    "",
    "T - Tone（语气）",
    "  • 正式、友好、鼓励性...",
    "",
    "A - Audience（受众）",
    "  • 写给谁看，他们的知识水平",
    "",
    "R - Response（响应）",
    "  • 输出格式、长度、结构"
])

# 13. 章节：实战示例
add_section_slide(prs, "4. 实战示例")

# 14. 对比示例 1
add_comparison_slide(
    prs,
    "示例 1：写代码",
    "写个爬虫",
    """你是一位 Python 专家

任务：编写一个新闻爬虫
- 目标网站：example.com
- 功能：抓取标题、链接、发布时间
- 要求：
  • 使用 requests + BeautifulSoup
  • 添加异常处理
  • 数据保存为 JSON
  • 添加 User-Agent 避免被封

输出：完整可运行的代码，包含注释"""
)

# 15. 对比示例 2
add_comparison_slide(
    prs,
    "示例 2：写文档",
    "写个产品说明",
    """你是一位资深技术文档工程师

任务：为新的 API 接口编写使用文档

背景：这是一个用户认证 API，面向开发者

输出要求：
- 包含接口说明、请求参数、返回格式
- 提供 3 个代码示例（Python/Java/cURL）
- 常见错误码说明
- 使用 Markdown 格式
- 语言简洁专业"""
)

# 16. 对比示例 3
add_comparison_slide(
    prs,
    "示例 3：数据分析",
    "分析一下数据",
    """你是一位数据分析师

任务：分析 Q4 销售数据

数据：[附上 CSV/Excel 数据]

分析维度：
- 同比/环比增长率
- 各产品线贡献占比
- 区域销售分布
- Top 10 客户分析

输出：
- 关键发现（3-5 条）
- 可视化建议（图表类型）
- 改进建议（ actionable insights）"""
)

# 17. 章节：常见错误
add_section_slide(prs, "5. 常见错误与优化")

# 18. 常见错误
add_content_slide(prs, "常见错误", [
    "❌ 太模糊",
    "  • 「帮我写点什么」→ 写什么？给谁看？",
    "",
    "❌ 缺少上下文",
    "  • 「优化这段代码」→ 优化目标？性能？可读性？",
    "",
    "❌ 信息过载",
    "  • 一次性给太多信息，AI 抓不住重点",
    "",
    "❌ 没有输出要求",
    "  • AI 不知道你要什么格式",
    "",
    "❌ 忽略角色设定",
    "  • 不同角色输出质量差异巨大"
])

# 19. 优化技巧
add_content_slide(prs, "优化技巧", [
    "✅ 添加角色",
    "  • 「你是一位...」",
    "",
    "✅ 分步骤说明",
    "  • 「第一步...第二步...」",
    "",
    "✅ 提供示例",
    "  • 「参考这个格式：...」",
    "",
    "✅ 指定输出格式",
    "  • 「用表格/列表/JSON 输出」",
    "",
    "✅ 设置约束",
    "  • 「不超过 500 字」「包含 3 个要点」"
])

# 20. 迭代优化
add_code_slide(
    prs,
    "迭代优化示例",
    """第 1 版：「写个邮件」
↓
第 2 版：「写个请假邮件」
↓
第 3 版：「写个 3 天的病假邮件给老板」
↓
第 4 版：
「你是一位职场人士
任务：写请假邮件给直属领导
情况：突发高烧，需要请假 3 天
要求：
- 语气专业且诚恳
- 说明工作交接安排
- 表达尽快恢复的意愿
- 200 字以内」""",
    "好的 Prompt 是迭代出来的！"
)

# 21. 章节：高级技巧
add_section_slide(prs, "6. 高级技巧")

# 22. 思维链
add_content_slide(prs, "思维链（Chain of Thought）", [
    "让 AI 展示思考过程，提升准确性",
    "",
    "关键句式：",
    "  • 「请逐步思考...」",
    "  • 「让我们一步步分析...」",
    "  • 「先...然后...最后...」",
    "",
    "示例：",
    "  「请逐步分析这个问题：",
    "   1. 先理解问题本质",
    "   2. 列出可能的解决方案",
    "   3. 评估每个方案的优缺点",
    "   4. 给出最终建议」"
])

# 23. 少样本学习
add_content_slide(prs, "少样本学习（Few-Shot）", [
    "提供示例，让 AI 模仿",
    "",
    "示例：",
    "  「把下面句子翻译成英文：",
    "  ",
    "  例子 1：",
    "  输入：今天天气真好",
    "  输出：The weather is really nice today",
    "  ",
    "  例子 2：",
    "  输入：我喜欢吃苹果",
    "  输出：I like to eat apples",
    "  ",
    "  现在翻译：这个电影很有趣」"
])

# 24. 自我验证
add_content_slide(prs, "自我验证（Self-Verification）", [
    "让 AI 检查自己的输出",
    "",
    "关键句式：",
    "  • 「完成后请检查是否有错误」",
    "  • 「请验证这个方案是否可行」",
    "  • 「如果有遗漏，请补充」",
    "",
    "示例：",
    "  「请编写一个函数，完成后：",
    "  1. 检查边界条件",
    "  2. 验证时间复杂度",
    "  3. 提供测试用例」"
])

# 25. 章节：总结
add_section_slide(prs, "7. 总结与实践")

# 26. 总结
add_content_slide(prs, "核心要点总结", [
    "1️⃣ Prompt 决定 AI 输出质量",
    "",
    "2️⃣ 好 Prompt = 角色 + 任务 + 上下文 + 约束 + 格式",
    "",
    "3️⃣ 使用框架：RTF / CO-STAR",
    "",
    "4️⃣ 避免模糊，要具体明确",
    "",
    "5️⃣ 迭代优化，逐步完善",
    "",
    "6️⃣ 高级技巧：思维链、少样本、自我验证"
])

# 27. 实践建议
add_content_slide(prs, "实践建议", [
    "🎯 建立自己的 Prompt 库",
    "  • 收集好用的 Prompt 模板",
    "",
    "📝 记录优化过程",
    "  • 什么改动让效果更好",
    "",
    "🔄 持续迭代",
    "  • 根据输出调整 Prompt",
    "",
    "💡 多尝试不同框架",
    "  • 找到最适合你的方式",
    "",
    "📚 学习优秀案例",
    "  • PromptBase、GitHub 等社区"
])

# 28. 结束页
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 背景
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
shape.line.fill.background()

# 文字
txBox = slide.shapes.add_textbox(
    Inches(0.5), Inches(2),
    prs.slide_width - Inches(1), Inches(3)
)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Q & A\n\n谢谢观看！"
p.font.size = Pt(54)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(255, 255, 255)

# 保存
output_path = '/Users/sunyu.maner/.openclaw/workspace/prompt_design_guide.pptx'
prs.save(output_path)

print(f"✅ PPT 创建完成：{output_path}")
print(f"📊 共 {len(prs.slides)} 页")
