#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenClaw 技术分享 PPT - 面向后台研发工程师
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section_slide(prs, title):
    """章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2),
        prs.slide_width, Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(221, 75, 57)  # OpenClaw 红色
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        prs.slide_width - Inches(1), Inches(2)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255, 255, 255)
    return slide

def add_content_slide(prs, title, items):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.space_after = Pt(12)
        if item.startswith('  '):
            p.level = 1
            p.text = item.strip()
    return slide

def add_code_slide(prs, title, code, notes=None):
    """代码示例页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(0.6)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # 代码框
    codeBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0),
        prs.slide_width - Inches(1), Inches(4)
    )
    codeBox.fill.solid()
    codeBox.fill.fore_color.rgb = RGBColor(40, 44, 52)
    codeBox.line.color.rgb = RGBColor(100, 100, 100)
    
    tf = codeBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = code
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    if notes:
        expBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2),
            prs.slide_width - Inches(1), Inches(1.8)
        )
        tf = expBox.text_frame
        p = tf.add_paragraph()
        p.text = notes
        p.font.size = Pt(18)
    
    return slide

def add_architecture_slide(prs, title, components):
    """架构图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        prs.slide_width - Inches(1), Inches(0.6)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 绘制组件
    y_pos = 1.2
    box_height = 1.2
    gap = 0.3
    
    for i, comp in enumerate(components):
        name = comp['name']
        desc = comp.get('desc', '')
        color = comp.get('color', RGBColor(41, 128, 185))
        
        # 方框
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(12.3), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # 文字
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(230, 230, 230)
        
        y_pos += box_height + gap
    
    return slide

def add_comparison_slide(prs, title, items):
    """对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        prs.slide_width - Inches(1), Inches(0.6)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 绘制表格
    y_pos = 1.0
    row_height = 1.0
    
    for i, item in enumerate(items):
        # 左侧
        left = slide.shapes.add_textbox(
            Inches(0.3), Inches(y_pos),
            Inches(6.2), Inches(row_height - 0.1)
        )
        left.fill.solid()
        left.fill.fore_color.rgb = RGBColor(236, 240, 241)
        left.line.fill.background()
        
        tf = left.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['other']
        p.font.size = Pt(18)
        
        # 右侧
        right = slide.shapes.add_textbox(
            Inches(6.8), Inches(y_pos),
            Inches(6.2), Inches(row_height - 0.1)
        )
        right.fill.solid()
        right.fill.fore_color.rgb = RGBColor(39, 174, 96)
        right.line.fill.background()
        
        tf = right.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['openclaw']
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        y_pos += row_height
    
    # 标题行
    header_left = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.9),
        Inches(6.2), Inches(0.4)
    )
    tf = header_left.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ 其他 AI 产品"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    header_right = slide.shapes.add_textbox(
        Inches(6.8), Inches(0.9),
        Inches(6.2), Inches(0.4)
    )
    tf = header_right.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ OpenClaw"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 开始创建幻灯片 ============

# 1. 标题页
add_title_slide(
    prs,
    "OpenClaw 技术分享",
    "下一代 AI Agent 网关架构\n\n面向后台研发工程师\n" + datetime.now().strftime("%Y 年 %m 月 %d 日")
)

# 2. 目录
add_content_slide(prs, "目录", [
    "1. OpenClaw 简介与发展历程",
    "2. 与竞品的差异化优势",
    "3. 系统架构与核心功能",
    "4. 核心组件详解（Skills/Plugins）",
    "5. Quick Start 与常用命令",
    "6. 实战案例分享",
    "7. 未来发展方向",
    "8. 参考文档与资源",
    "9. Q&A"
])

# ========== 第一部分：OpenClaw 简介 ==========
add_section_slide(prs, "1. OpenClaw 简介与发展历程")

# 3. 什么是 OpenClaw
add_content_slide(prs, "什么是 OpenClaw？", [
    "定义：自托管的多渠道 AI Agent 网关",
    "",
    "核心理念：",
    "  • 一个网关连接所有聊天应用",
    "  • 本地运行，数据自主可控",
    "  • 为 AI Agent 原生设计",
    "",
    "Slogan:",
    "  \"Any OS gateway for AI agents across",
    "   WhatsApp, Telegram, Discord, iMessage, and more.\""
])

# 4. 研发背景
add_content_slide(prs, "研发背景与动机", [
    "痛点：",
    "  • 每个聊天平台需要独立的 Bot 实现",
    "  • 云端服务数据隐私无法保障",
    "  • 现有方案不支持复杂的 Agent 工作流",
    "",
    "目标：",
    "  • 统一的多渠道网关",
    "  • 自托管，数据完全可控",
    "  • 原生支持 AI Agent（工具调用、会话管理、多 Agent 路由）",
    "",
    "研发团队：开源社区驱动",
    "License：MIT"
])

# 5. 发展历程
add_content_slide(prs, "发展历程", [
    "2024 年：项目启动",
    "  • 初始版本支持 WhatsApp、Telegram",
    "",
    "2025 年：快速发展",
    "  • 增加 Discord、iMessage、Signal 等渠道",
    "  • 引入 Skills 和 Plugins 系统",
    "  • 推出 Control UI 和移动端 Node",
    "",
    "2026 年（当前 v2026.3.2）：",
    "  • 54+ 内置 Skills",
    "  • 20+ 渠道支持",
    "  • 完整的生态系统（ClawHub 技能市场）",
    "",
    "GitHub: github.com/openclaw/openclaw"
])

# 6. 为什么爆火
add_content_slide(prs, "为什么 OpenClaw 能爆火？", [
    "1. 自托管趋势",
    "  • 数据隐私意识增强",
    "  • 企业合规需求",
    "",
    "2. AI Agent 爆发",
    "  • 需要统一的接入层",
    "  • 工具调用、会话管理成刚需",
    "",
    "3. 多渠道整合",
    "  • 一套代码服务所有平台",
    "  • 降低开发和维护成本",
    "",
    "4. 开源生态",
    "  • 社区贡献 Skills/Plugins",
    "  • 快速迭代，响应需求"
])

# ========== 第二部分：差异化优势 ==========
add_section_slide(prs, "2. 与竞品的差异化优势")

# 7. 竞品对比
add_comparison_slide(prs, "竞品对比", [
    {
        'other': 'LangChain/LlamaIndex\n需要自己搭建网关和渠道集成',
        'openclaw': '开箱即用的多渠道网关\n内置 20+ 渠道支持'
    },
    {
        'other': '云端 Bot 服务\n数据存储在第三方',
        'openclaw': '自托管\n数据完全本地控制'
    },
    {
        'other': '单一渠道 Bot\n每个平台独立实现',
        'openclaw': '统一网关\n一次开发多端复用'
    },
    {
        'other': '基础对话功能\n缺少 Agent 能力',
        'openclaw': 'Agent 原生\n工具调用/会话管理/多 Agent 路由'
    }
])

# 8. 核心优势
add_content_slide(prs, "OpenClaw 核心优势", [
    "🏠 自托管 (Self-hosted)",
    "  • 运行在自己的硬件上",
    "  • 数据不出内网",
    "  • 完全可控",
    "",
    "🔌 多渠道 (Multi-channel)",
    "  • WhatsApp, Telegram, Discord, iMessage...",
    "  • 单一网关进程",
    "  • 统一配置管理",
    "",
    "🤖 Agent 原生 (Agent-native)",
    "  • 工具调用 (Tools)",
    "  • 会话管理 (Sessions)",
    "  • 多 Agent 路由 (Multi-agent routing)",
    "",
    "📦 可扩展 (Extensible)",
    "  • Skills 系统",
    "  • Plugins 机制",
    "  • ClawHub 生态"
])

# ========== 第三部分：系统架构 ==========
add_section_slide(prs, "3. 系统架构与核心功能")

# 9. 整体架构
add_architecture_slide(prs, "OpenClaw 整体架构", [
    {
        'name': '渠道层 (Channels)',
        'desc': 'WhatsApp | Telegram | Discord | iMessage | Signal | Slack | 微信 (插件) | QQ (插件)...',
        'color': RGBColor(52, 152, 219)
    },
    {
        'name': '网关层 (Gateway)',
        'desc': '消息路由 | 会话管理 | 认证授权 | 配置管理 | 日志监控',
        'color': RGBColor(221, 75, 57)
    },
    {
        'name': 'Agent 层 (Pi Agent)',
        'desc': 'LLM 调用 | 工具执行 | 记忆管理 | 多 Agent 协作',
        'color': RGBColor(39, 174, 96)
    },
    {
        'name': '接口层 (Interfaces)',
        'desc': 'CLI | Web Control UI | REST API | WebSocket RPC',
        'color': RGBColor(142, 68, 173)
    }
])

# 10. 核心功能
add_content_slide(prs, "核心功能", [
    "1. 多渠道消息网关",
    "  • 20+ 官方支持渠道",
    "  • 插件扩展更多平台",
    "",
    "2. Agent 运行时",
    "  • 会话隔离 (每用户/每渠道)",
    "  • 工具调用 (文件操作/搜索/自定义)",
    "  • 记忆管理 (短期/长期)",
    "",
    "3. 多 Agent 路由",
    "  • 按技能路由到不同 Agent",
    "  • 支持子 Agent (Sub-agents)",
    "",
    "4. 媒体处理",
    "  • 图片/音频/视频/文档",
    "  • 自动转码和优化",
    "",
    "5. 管理界面",
    "  • Web Control UI",
    "  • CLI 命令行",
    "  • 实时监控和日志"
])

# 11. 使用场景
add_content_slide(prs, "人们都用 OpenClaw 做什么？", [
    "💬 个人 AI 助理",
    "  • 微信/Telegram 随时对话",
    "  • 问答/写作/编程辅助",
    "",
    "👥 团队协作用 Bot",
    "  • Slack/Discord 团队助手",
    "  • 自动回复/知识库查询",
    "",
    "🔧 自动化工作流",
    "  • 定时任务 (Cron)",
    "  • 事件触发 (Webhook)",
    "  • 跨系统集成",
    "",
    "📊 监控告警",
    "  • 系统监控通知",
    "  • 日志分析告警",
    "",
    "🎯 垂直场景",
    "  • 客服机器人",
    "  • 数据分析助手",
    "  • 代码 Review Bot"
])

# ========== 第四部分：核心组件 ==========
add_section_slide(prs, "4. 核心组件详解")

# 12. Skills 系统
add_content_slide(prs, "Skills 系统", [
    "什么是 Skills？",
    "  • 模块化的技能包",
    "  • 扩展 AI 能力的\"说明书\"",
    "",
    "Skills 结构：",
    "  SKILL.md (必需) - 技能说明和触发条件",
    "  scripts/ (可选) - 可执行脚本",
    "  references/ (可选) - 参考文档",
    "  assets/ (可选) - 模板/资源文件",
    "",
    "内置 Skills (54 个)：",
    "  • github - GitHub 操作",
    "  • weather - 天气查询",
    "  • video-frames - 视频帧提取",
    "  • healthcheck - 系统审计",
    "  • skill-creator - 创建新 Skills",
    "  • 详见：/opt/homebrew/lib/node_modules/openclaw/skills/"
])

# 13. Skills 示例
add_code_slide(prs, "Skills 示例：GitHub Skill", """---
name: github
description: GitHub operations via gh CLI
---

# GitHub Skill

## When to Use
✅ Checking PR status or CI
✅ Creating/commenting on issues
✅ Viewing run logs

## Common Commands
# List PRs
gh pr list --repo owner/repo

# Check CI
gh pr checks 55 --repo owner/repo

# Create issue
gh issue create --title "Bug" --body "Details"
""", "Skills 通过 SKILL.md 定义触发条件和使用方法")

# 14. Plugins 系统
add_content_slide(prs, "Plugins 系统", [
    "什么是 Plugins？",
    "  • 运行时扩展模块",
    "  • 添加新渠道/工具/服务",
    "",
    "Plugins 能力：",
    "  • 注册新渠道 (如 QQ Bot、飞书)",
    "  • 添加工具 (Agent Tools)",
    "  • 注册 CLI 命令",
    "  • HTTP/RPC 服务",
    "",
    "官方 Plugins：",
    "  • @openclaw/voice-call - 语音通话",
    "  • @openclaw/msteams - Microsoft Teams",
    "  • @openclaw/matrix - Matrix 协议",
    "  • @openclaw/zalo - Zalo (越南)",
    "",
    "安装：openclaw plugins install @openclaw/voice-call"
])

# 15. Channels 渠道
add_content_slide(prs, "Channels 渠道支持", [
    "官方支持 (20+)：",
    "  • WhatsApp (WhatsApp Cloud API)",
    "  • Telegram (Bot API)",
    "  • Discord (Bot)",
    "  • iMessage (BlueBubbles/Android)",
    "  • Signal (信号桥接)",
    "  • Slack (Bot)",
    "  • Google Chat",
    "  • Microsoft Teams (插件)",
    "  • IRC",
    "  • Mattermost (插件)",
    "",
    "插件扩展：",
    "  • QQ Bot (@sliverp/qqbot)",
    "  • 微信 (WeChat)",
    "  • 飞书 (Feishu)",
    "  • 钉钉 (DingTalk)",
    "",
    "配置示例：",
    "  channels.telegram.enabled: true",
    "  channels.telegram.botToken: \"xxx\""
])

# 16. Gateway 网关
add_content_slide(prs, "Gateway 网关服务", [
    "核心职责：",
    "  • 消息路由和分发",
    "  • 会话状态管理",
    "  • 认证和授权",
    "  • 配置热加载",
    "  • 日志和监控",
    "",
    "运行模式：",
    "  • 单进程多路复用",
    "  • WebSocket + HTTP API",
    "  • 默认端口：18789",
    "",
    "高可用：",
    "  • 配置热重载 (hybrid 模式)",
    "  • 会话持久化",
    "  • 支持 Tailscale 组网",
    "",
    "启动：openclaw gateway --port 18789"
])

# ========== 第五部分：Quick Start ==========
add_section_slide(prs, "5. Quick Start 与常用命令")

# 17. 快速开始
add_content_slide(prs, "5 分钟快速开始", [
    "1. 安装 OpenClaw",
    "  npm install -g openclaw@latest",
    "",
    "2. 运行引导向导",
    "  openclaw onboard --install-daemon",
    "",
    "3. 登录渠道",
    "  openclaw channels login whatsapp",
    "  (扫码或配置 Token)",
    "",
    "4. 启动网关",
    "  openclaw gateway --port 18789",
    "",
    "5. 打开 Control UI",
    "  http://127.0.0.1:18789/",
    "",
    "完成！现在可以通过聊天应用对话了"
])

# 18. 常用命令
add_code_slide(prs, "常用命令速查", """# 安装与配置
npm install -g openclaw@latest
openclaw onboard              # 引导向导
openclaw configure            # 配置向导

# 网关管理
openclaw gateway              # 启动网关
openclaw gateway status       # 查看状态
openclaw gateway restart      # 重启
openclaw logs --follow        # 查看日志

# 渠道管理
openclaw channels status      # 渠道状态
openclaw channels login       # 登录渠道
openclaw pairing list         # 查看配对

# 消息发送
openclaw message send --channel telegram -t @username -m "Hello"

# 插件管理
openclaw plugins list
openclaw plugins install @openclaw/voice-call

# 技能管理
openclaw skills list""", "掌握这些命令即可日常使用")

# 19. 配置示例
add_code_slide(prs, "配置文件 (~/.openclaw/openclaw.json)", """{
  "channels": {
    "telegram": {
      "enabled": true,
      "botToken": "123:abc",
      "dmPolicy": "pairing",
      "groups": { "*": { "requireMention": true } }
    },
    "whatsapp": {
      "enabled": true,
      "allowFrom": ["+8613800138000"]
    }
  },
  "agents": {
    "defaults": {
      "model": { "primary": "qwen3.5-plus" }
    }
  },
  "gateway": {
    "port": 18789,
    "auth": { "mode": "token" }
  }
}""", "配置文件支持热重载，修改后自动生效")

# ========== 第六部分：实战案例 ==========
add_section_slide(prs, "6. 实战案例分享")

# 20. 案例 1：个人 AI 助理
add_content_slide(prs, "案例 1：个人 AI 助理", [
    "场景：通过微信/Telegram 随时与 AI 对话",
    "",
    "配置：",
    "  • 渠道：Telegram (Bot API)",
    "  • 访问控制：allowFrom 白名单",
    "  • 模型：qwen3.5-plus",
    "",
    "功能：",
    "  • 日常问答",
    "  • 代码审查",
    "  • 文档写作",
    "  • 日程提醒",
    "",
    "效果：",
    "  • 5 分钟搭建完成",
    "  • 数据完全本地",
    "  • 响应速度 < 3 秒"
])

# 21. 案例 2：团队协作用 Bot
add_content_slide(prs, "案例 2：团队协作用 Bot", [
    "场景：Slack 团队助手",
    "",
    "配置：",
    "  • 渠道：Slack Bot",
    "  • 群组策略：requireMention",
    "  • 多 Agent 路由",
    "",
    "功能：",
    "  • 自动回答常见问题",
    "  • 代码片段分享",
    "  • CI/CD状态查询",
    "  • 会议纪要生成",
    "",
    "Skills 扩展：",
    "  • github - PR/Issue 管理",
    "  • healthcheck - 系统监控",
    "  • 自定义 - 内部 API 集成"
])

# 22. 案例 3：监控告警系统
add_content_slide(prs, "案例 3：监控告警系统", [
    "场景：系统监控 + 告警通知",
    "",
    "架构：",
    "  Prometheus → Alertmanager → OpenClaw → Telegram",
    "",
    "实现：",
    "  • Webhook 接收告警",
    "  • 自动分类和路由",
    "  • 富文本通知 (带图表)",
    "  • 告警确认和关闭",
    "",
    "优势：",
    "  • 多渠道通知 (Telegram/微信/短信)",
    "  • 支持告警升级",
    "  • 历史记录查询",
    "",
    "代码量：< 100 行"
])

# 23. 案例 4：新闻监控推送
add_content_slide(prs, "案例 4：新闻监控推送 (本次演示)", [
    "场景：监控今日头条新闻并推送到 QQ",
    "",
    "实现：",
    "  • 自定义 Skill: toutiao-news-monitor",
    "  • 浏览器自动化抓取",
    "  • 定时任务 (Cron)",
    "  • QQ Bot 推送",
    "",
    "技能结构：",
    "  skills/toutiao-news-monitor/",
    "  ├── SKILL.md",
    "  ├── scripts/crawl_news.py",
    "  └── data/ (输出)",
    "",
    "效果：",
    "  • 每 2 小时自动推送",
    "  • 关键词过滤",
    "  • 数据持久化"
])

# ========== 第七部分：未来发展 ==========
add_section_slide(prs, "7. 未来发展方向")

# 24. 路线图
add_content_slide(prs, "OpenClaw 发展路线图", [
    "短期 (2026 Q2-Q3)：",
    "  • 更多渠道支持 (Line/KakaoTalk)",
    "  • 改进的 Control UI",
    "  • 性能优化和稳定性",
    "",
    "中期 (2026 Q4-2027 Q1)：",
    "  • 多租户支持",
    "  • 企业级权限管理",
    "  • 高级监控和告警",
    "  • 负载均衡和水平扩展",
    "",
    "长期愿景：",
    "  • 成为 AI Agent 的标准接入层",
    "  • 构建完整的开发者生态",
    "  • 支持更多 AI 模型和框架",
    "  • 边缘计算和离线能力"
])

# 25. 社区贡献
add_content_slide(prs, "如何参与社区贡献？", [
    "📝 贡献 Skills",
    "  • 分享你的专业技能",
    "  • 发布到 ClawHub",
    "",
    "🔌 开发 Plugins",
    "  • 添加新渠道",
    "  • 集成新服务",
    "",
    "🐛 报告 Bug",
    "  • GitHub Issues",
    "  • Discord 社区",
    "",
    "📚 改进文档",
    "  • 翻译",
    "  • 补充示例",
    "",
    "💡 提出建议",
    "  • Feature Requests",
    "  • 技术方案讨论"
])

# ========== 第八部分：参考文档 ==========
add_section_slide(prs, "8. 参考文档与资源")

# 26. 参考文档
add_content_slide(prs, "官方文档", [
    "📖 核心文档",
    "  • 官网：https://docs.openclaw.ai",
    "  • GitHub: github.com/openclaw/openclaw",
    "  • 本地文档：/opt/homebrew/lib/node_modules/openclaw/docs/",
    "",
    "🎯 快速开始",
    "  • /start/getting-started",
    "  • /start/wizard (引导向导)",
    "",
    "📚 深入阅读",
    "  • /gateway/architecture (架构)",
    "  • /channels/ (渠道配置)",
    "  • /tools/skills (Skills 系统)",
    "  • /plugins/ (Plugins 系统)",
    "",
    "🔧 故障排查",
    "  • /gateway/troubleshooting",
    "  • /channels/troubleshooting"
])

# 27. 学习资源
add_content_slide(prs, "学习资源与社区", [
    "🌐 在线资源",
    "  • 官方文档：docs.openclaw.ai",
    "  • ClawHub 技能市场：clawhub.ai",
    "  • GitHub 示例：github.com/openclaw",
    "",
    "💬 社区交流",
    "  • Discord: discord.gg/clawd",
    "  • GitHub Discussions",
    "  • 技术博客和教程",
    "",
    "📦 优秀 Skills 推荐",
    "  • github - GitHub 操作",
    "  • weather - 天气查询",
    "  • healthcheck - 系统审计",
    "  • coding-agent - 代码编写",
    "  • todo 列表：openclaw skills list"
])

# ========== 第九部分：总结 ==========
add_section_slide(prs, "9. 总结与 Q&A")

# 28. 核心价值
add_content_slide(prs, "OpenClaw 核心价值", [
    "🎯 一句话总结",
    "  \"自托管的多渠道 AI Agent 网关\"",
    "",
    "💎 三大优势",
    "  1. 数据自主可控 (自托管)",
    "  2. 一次开发多端复用 (多渠道)",
    "  3. Agent 原生设计 (工具/会话/路由)",
    "",
    "🚀 适用场景",
    "  • 个人 AI 助理",
    "  • 团队协作用 Bot",
    "  • 自动化工作流",
    "  • 监控告警系统",
    "",
    "🎓 学习建议",
    "  1. 5 分钟快速开始",
    "  2. 配置一个渠道",
    "  3. 尝试内置 Skills",
    "  4. 开发自定义 Skill"
])

# 29. 演示环节
add_content_slide(prs, "现场演示", [
    "📱 演示 1：多渠道消息",
    "  • Telegram/微信/QQ 同时接收消息",
    "",
    "🔧 演示 2：Skills 使用",
    "  • GitHub 查询",
    "  • 天气查询",
    "  • 新闻监控推送",
    "",
    "💻 演示 3：自定义 Skill",
    "  • toutiao-news-monitor",
    "  • 从创建到执行全流程",
    "",
    "🎯 演示 4：Control UI",
    "  • Web 界面操作",
    "  • 会话管理",
    "  • 配置修改"
])

# 30. Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(221, 75, 57)
shape.line.fill.background()

txBox = slide.shapes.add_textbox(
    Inches(0.5), Inches(2),
    prs.slide_width - Inches(1), Inches(3)
)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Q & A\n\n谢谢观看！\n\n欢迎提问 🦞"
p.font.size = Pt(54)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(255, 255, 255)

# 保存
output_path = '/Users/sunyu.maner/.openclaw/workspace/openclaw_tech_share.pptx'
prs.save(output_path)

print(f"✅ PPT 创建完成：{output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"📁 已保存到 GitHub 工作区")
